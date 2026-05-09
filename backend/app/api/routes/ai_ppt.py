from __future__ import annotations

import json
import re
import uuid
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pydantic import BaseModel

from app.api.routes.auth import AccountInfo, get_current_account

router = APIRouter()
MAX_PPT_UPLOAD_BYTES = 50 * 1024 * 1024


class AIPPTFeishuActionRequest(BaseModel):
    filename: str
    action: str
    target: str = ""
    folder_token: str = ""
    wiki_token: str = ""
    message: str = ""
    session_id: str = ""


def ai_ppt_base_dir() -> Path:
    path = Path(__file__).resolve().parents[2] / "data" / "ai_ppt"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _backend_dir() -> Path:
    return Path(__file__).resolve().parents[3]


def _template_dir() -> Path:
    path = ai_ppt_base_dir() / "templates"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _builtin_template_dir() -> Path:
    return Path(__file__).resolve().parents[2] / "assets" / "ai_ppt_templates"


def _preview_dir() -> Path:
    path = ai_ppt_base_dir() / "previews"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _source_dir() -> Path:
    path = ai_ppt_base_dir() / "sources"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _template_meta_path() -> Path:
    return _template_dir() / "templates.json"


def _load_template_meta() -> list[dict]:
    path = _template_meta_path()
    if not path.exists():
        return []
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return []
    return data if isinstance(data, list) else []


def _save_template_meta(items: list[dict]) -> None:
    _template_meta_path().write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")


def _safe_name(name: str) -> str:
    stem = Path(name or "presentation.pptx").stem
    cleaned = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", stem).strip("_")
    return cleaned[:40] or "presentation"


def _is_relative_to(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
        return True
    except ValueError:
        return False


def _default_templates() -> list[dict]:
    return [
        {
            "id": "ppt_master_mckinsey_customer_loyalty",
            "name": "咨询汇报",
            "style": "咨询",
            "description": "参考 ppt-master 示例模板，适合战略、经营分析和管理汇报。",
            "builtin": True,
            "accent": "#111827",
            "filename": "ppt_master_mckinsey_customer_loyalty.pptx",
            "download_url": "/api/v1/ai-ppt/files/ppt_master_mckinsey_customer_loyalty.pptx",
        },
        {
            "id": "ppt_master_google_annual_report",
            "name": "年度报告",
            "style": "商务",
            "description": "参考 ppt-master 年报模板，适合复盘、数据总结和组织汇报。",
            "builtin": True,
            "accent": "#4285f4",
            "filename": "ppt_master_google_annual_report.pptx",
            "download_url": "/api/v1/ai-ppt/files/ppt_master_google_annual_report.pptx",
        },
        {
            "id": "ppt_master_dark_tech",
            "name": "暗色科技",
            "style": "科技",
            "description": "参考 ppt-master 深色科技模板，适合 AI 产品、技术方案和发布会。",
            "builtin": True,
            "accent": "#0f172a",
            "filename": "ppt_master_dark_tech.pptx",
            "download_url": "/api/v1/ai-ppt/files/ppt_master_dark_tech.pptx",
        },
    ]


def resolve_ai_ppt_file(file_id: Optional[str]) -> Optional[Path]:
    if not file_id:
        return None
    candidates = [
        _source_dir() / f"{file_id}.pptx",
        _template_dir() / f"{file_id}.pptx",
        _builtin_template_dir() / f"{file_id}.pptx",
        ai_ppt_base_dir() / file_id,
        _source_dir() / file_id,
        _template_dir() / file_id,
        _builtin_template_dir() / file_id,
    ]
    roots = [ai_ppt_base_dir().resolve(), _builtin_template_dir().resolve()]
    for candidate in candidates:
        target = candidate.resolve()
        if target.is_file() and target.suffix.lower() == ".pptx" and any(_is_relative_to(target, root) for root in roots):
            return target
    return None


def resolve_ai_ppt_template(template_id: Optional[str]) -> Optional[Path]:
    return resolve_ai_ppt_file(template_id)


@router.get("/ai-ppt/templates")
async def list_ai_ppt_templates() -> dict:
    return {"code": 0, "data": {"templates": _default_templates() + _load_template_meta()}}


@router.post("/ai-ppt/source")
async def upload_ai_ppt_source(
    file: UploadFile = File(...),
    account: AccountInfo = Depends(get_current_account),
) -> dict:
    return await _save_uploaded_ppt(file=file, owner=account.account, kind="source")


@router.post("/ai-ppt/templates")
async def upload_ai_ppt_template(
    name: str = Form(default=""),
    style: str = Form(default="商务"),
    description: str = Form(default=""),
    file: UploadFile = File(...),
    account: AccountInfo = Depends(get_current_account),
) -> dict:
    saved = await _save_uploaded_ppt(file=file, owner=account.account, kind="template")
    item = {
        "id": saved["data"]["file_id"],
        "name": name.strip() or _safe_name(file.filename or "模板"),
        "style": style.strip() or "商务",
        "description": description.strip() or "用户上传模板",
        "filename": saved["data"]["stored_filename"],
        "download_url": saved["data"]["download_url"],
        "builtin": False,
        "accent": "#2454d6",
        "created_at": datetime.now().isoformat(),
    }
    items = [template for template in _load_template_meta() if template.get("id") != item["id"]]
    items.append(item)
    _save_template_meta(items)
    return {"code": 0, "data": item}


async def _save_uploaded_ppt(file: UploadFile, owner: str, kind: str) -> dict:
    if not (file.filename or "").lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="只支持 .pptx 文件")
    content = await file.read()
    if len(content) > MAX_PPT_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail="PPT 文件不能超过 50MB")

    file_id = uuid.uuid4().hex
    stored_filename = f"{file_id}.pptx"
    target_dir = _template_dir() if kind == "template" else _source_dir()
    target = target_dir / stored_filename
    target.write_bytes(content)
    return {
        "code": 0,
        "data": {
            "file_id": file_id,
            "filename": file.filename,
            "stored_filename": stored_filename,
            "size": len(content),
            "owner": owner,
            "kind": kind,
            "download_url": f"/api/v1/ai-ppt/files/{stored_filename}",
        },
    }


@router.get("/ai-ppt/files/{filename}")
async def download_ai_ppt(filename: str) -> FileResponse:
    target = _resolve_download_file(filename)
    if not target:
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(
        target,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=target.name,
    )


@router.get("/ai-ppt/files/{filename}/preview")
async def preview_ai_ppt(filename: str) -> dict:
    target = _resolve_download_file(filename)
    if not target:
        raise HTTPException(status_code=404, detail="File not found")
    return {"code": 0, "data": _render_preview(target)}


@router.get("/ai-ppt/previews/{preview_id}/{image_name}")
async def get_ai_ppt_preview_image(preview_id: str, image_name: str) -> FileResponse:
    target = (_preview_dir() / preview_id / image_name).resolve()
    if not _is_relative_to(target, _preview_dir().resolve()):
        raise HTTPException(status_code=404, detail="Preview not found")
    if not target.is_file() or target.suffix.lower() != ".png":
        raise HTTPException(status_code=404, detail="Preview not found")
    return FileResponse(target, media_type="image/png")


@router.post("/ai-ppt/actions")
async def run_ai_ppt_feishu_action(
    request: AIPPTFeishuActionRequest,
    account: AccountInfo = Depends(get_current_account),
) -> dict:
    target = _resolve_download_file(request.filename)
    if not target:
        raise HTTPException(status_code=404, detail="PPT file not found")
    from app.skills.lark_cli.skill import LarkCLISkill

    result = await _run_feishu_action(LarkCLISkill(), request, target, account.account)
    return {"code": 0, "data": result}


async def _run_feishu_action(skill: Any, request: AIPPTFeishuActionRequest, target_path: Path, user_id: str) -> dict:
    action = request.action.strip()
    relative_file = _ppt_cli_relative_path(target_path)
    executed: list[dict[str, Any]] = []

    if action == "upload":
        command = _build_upload_command(request, relative_file)
        success, stdout, stderr = await skill.execute_command(command, 60, user_id)
        executed.append(_command_result(command, success, stdout, stderr, "上传 PPTX 到飞书云文档"))
        return _action_response(skill, success, stdout, stderr, executed)

    if action == "send_group":
        if not request.target.strip():
            raise HTTPException(status_code=400, detail="请填写群聊名称")
        search_command = f"lark-cli im +chat-search --query {_quote_cli_arg(request.target.strip())} --format json --as user"
        success, stdout, stderr = await skill.execute_command(search_command, 30, user_id)
        executed.append(_command_result(search_command, success, stdout, stderr, "搜索目标群聊"))
        if not success:
            return _action_response(skill, False, stdout, stderr, executed)
        chat_id = _extract_first_value(stdout, ("chat_id",), "oc_")
        if not chat_id:
            return _action_response(skill, False, stdout, "未能从群搜索结果中解析 chat_id", executed)
        send_command = (
            f"lark-cli im +messages-send --chat-id {_quote_cli_arg(chat_id)} "
            f"--file {_quote_cli_arg(relative_file)} --as user"
        )
        success, stdout, stderr = await skill.execute_command(send_command, 60, user_id)
        executed.append(_command_result(send_command, success, stdout, stderr, "发送 PPTX 到群聊"))
        return _action_response(skill, success, stdout, stderr, executed)

    if action == "send_person":
        if not request.target.strip():
            raise HTTPException(status_code=400, detail="请填写联系人姓名")
        search_command = f"lark-cli contact +search-user --query {_quote_cli_arg(request.target.strip())} --format json"
        success, stdout, stderr = await skill.execute_command(search_command, 30, user_id)
        executed.append(_command_result(search_command, success, stdout, stderr, "搜索联系人"))
        if not success:
            return _action_response(skill, False, stdout, stderr, executed)
        open_id = _extract_first_value(stdout, ("open_id", "user_id"), "ou_")
        if not open_id:
            return _action_response(skill, False, stdout, "未能从联系人搜索结果中解析 open_id", executed)
        send_command = (
            f"lark-cli im +messages-send --user-id {_quote_cli_arg(open_id)} "
            f"--file {_quote_cli_arg(relative_file)} --as user"
        )
        success, stdout, stderr = await skill.execute_command(send_command, 60, user_id)
        executed.append(_command_result(send_command, success, stdout, stderr, "发送 PPTX 给联系人"))
        return _action_response(skill, success, stdout, stderr, executed)

    raise HTTPException(status_code=400, detail="Unsupported PPT action")


def _resolve_download_file(filename: str) -> Optional[Path]:
    roots = [ai_ppt_base_dir().resolve(), _builtin_template_dir().resolve()]
    candidates = [
        (ai_ppt_base_dir() / filename).resolve(),
        (_source_dir() / filename).resolve(),
        (_template_dir() / filename).resolve(),
        (_builtin_template_dir() / filename).resolve(),
    ]
    for target in candidates:
        if target.is_file() and target.suffix.lower() == ".pptx" and any(_is_relative_to(target, root) for root in roots):
            return target
    return None


def _build_upload_command(request: AIPPTFeishuActionRequest, relative_file: str) -> str:
    command = f"lark-cli drive +upload --file {_quote_cli_arg(relative_file)} --as user"
    if request.folder_token.strip():
        command += f" --folder-token {_quote_cli_arg(request.folder_token.strip())}"
    if request.wiki_token.strip():
        command += f" --wiki-token {_quote_cli_arg(request.wiki_token.strip())}"
    return command


def _command_result(command: str, success: bool, stdout: str, stderr: str, reason: str) -> dict[str, Any]:
    return {
        "command": command,
        "reason": reason,
        "expected": "write",
        "success": success,
        "stdout": stdout.strip(),
        "stderr": stderr.strip(),
    }


def _action_response(skill: Any, success: bool, stdout: str, stderr: str, executed: list[dict[str, Any]]) -> dict:
    output = (stdout or stderr or "").strip()
    metadata: dict[str, Any] = {"executed_commands": executed}
    missing_scopes = skill._extract_missing_scopes(output) if hasattr(skill, "_extract_missing_scopes") else []
    if missing_scopes and hasattr(skill, "_build_scope_setup_metadata"):
        try:
            metadata.update(skill._build_scope_setup_metadata(None, missing_scopes))
        except Exception:
            scope_text = " ".join(missing_scopes)
            metadata.update(
                {
                    "setup_required": True,
                    "setup_scopes": missing_scopes,
                    "setup_steps": [
                        {
                            "key": "auth_login",
                            "title": "补充授权",
                            "command": f"lark-cli auth login --scope {_quote_cli_arg(scope_text)}",
                        }
                    ],
                    "setup_guide": "",
                }
            )
    if success:
        message = output or "飞书操作已完成"
    elif missing_scopes:
        message = f"飞书操作未完成：缺少权限 {', '.join(missing_scopes)}。请补充授权后回到 PPT 卡片重新执行。"
    else:
        message = output or "飞书操作失败"
    return {"success": success, "message": message, "metadata": metadata}


def _ppt_cli_relative_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(_backend_dir().resolve())).replace("\\", "/")
    except ValueError:
        return str(path)


def _extract_json_payload(text: str) -> Any:
    raw = (text or "").strip()
    if not raw:
        return None
    candidates = [raw]
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        candidates.append(raw[start : end + 1])
    for candidate in candidates:
        try:
            return json.loads(candidate)
        except Exception:
            continue
    return None


def _extract_first_value(text: str, keys: tuple[str, ...], prefix: str) -> str:
    payload = _extract_json_payload(text)
    seen: set[int] = set()

    def walk(value: Any) -> str:
        value_id = id(value)
        if value_id in seen:
            return ""
        seen.add(value_id)
        if isinstance(value, dict):
            for key in keys:
                item = str(value.get(key) or "")
                if item.startswith(prefix):
                    return item
            for item in value.values():
                found = walk(item)
                if found:
                    return found
        elif isinstance(value, list):
            for item in value:
                found = walk(item)
                if found:
                    return found
        return ""

    found = walk(payload)
    if found:
        return found
    match = re.search(rf"{re.escape(prefix)}[A-Za-z0-9_\\-]+", text or "")
    return match.group(0) if match else ""


def _quote_cli_arg(value: str) -> str:
    return '"' + value.replace('"', '\\"') + '"'


def _render_preview(ppt_path: Path) -> dict:
    prs = Presentation(str(ppt_path))
    preview_id = f"{ppt_path.stem}_{int(ppt_path.stat().st_mtime)}"
    out_dir = _preview_dir() / preview_id
    out_dir.mkdir(parents=True, exist_ok=True)

    width, height = 1280, 720
    scale_x = width / float(prs.slide_width)
    scale_y = height / float(prs.slide_height)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        image_name = f"slide_{index:03d}.png"
        image_path = out_dir / image_name
        texts: list[str] = []
        if not image_path.exists():
            image = Image.new("RGB", (width, height), "#ffffff")
            draw = ImageDraw.Draw(image)
            bg = _slide_background(slide)
            if bg:
                draw.rectangle([0, 0, width, height], fill=bg)
            for shape in slide.shapes:
                _draw_shape_preview(draw, image, shape, scale_x, scale_y, texts)
            image.save(image_path)
        else:
            texts = _slide_texts(slide)
        slides.append(
            {
                "index": index,
                "image_url": f"/api/v1/ai-ppt/previews/{preview_id}/{image_name}",
                "texts": [text.strip() for text in texts if text.strip()],
            }
        )
    return {"filename": ppt_path.name, "slide_count": len(slides), "slides": slides}


def _slide_texts(slide) -> list[str]:
    items = []
    for shape in slide.shapes:
        try:
            text = shape.text
        except Exception:
            text = ""
        if text:
            items.append(text)
    return items


def _slide_background(slide) -> Optional[str]:
    try:
        fill = slide.background.fill
        if fill.fore_color and fill.fore_color.rgb:
            return f"#{fill.fore_color.rgb}"
    except Exception:
        return None
    return None


def _draw_shape_preview(draw: ImageDraw.ImageDraw, image: Image.Image, shape, scale_x: float, scale_y: float, texts: list[str]) -> None:
    try:
        x = int(float(shape.left) * scale_x)
        y = int(float(shape.top) * scale_y)
        w = max(1, int(float(shape.width) * scale_x))
        h = max(1, int(float(shape.height) * scale_y))
    except Exception:
        return

    try:
        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
            picture = Image.open(BytesIO(shape.image.blob)).convert("RGBA")
            picture.thumbnail((w, h))
            image.paste(picture, (x, y), picture)
            return
    except Exception:
        pass

    fill = "#ffffff"
    outline = "#d0d5dd"
    try:
        if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
            fill = f"#{shape.fill.fore_color.rgb}"
    except Exception:
        fill = ""
    try:
        if shape.line and shape.line.color and shape.line.color.rgb:
            outline = f"#{shape.line.color.rgb}"
    except Exception:
        pass

    if fill:
        draw.rectangle([x, y, x + w, y + h], fill=fill, outline=outline)
    elif w > 20 and h > 20:
        draw.rectangle([x, y, x + w, y + h], outline="#e5e7eb")

    try:
        text = shape.text or ""
    except Exception:
        text = ""
    if text.strip():
        texts.append(text)
        _draw_wrapped_text(draw, text, x + 10, y + 8, max(20, w - 20), max(18, h - 16))


def _draw_wrapped_text(draw: ImageDraw.ImageDraw, text: str, x: int, y: int, width: int, height: int) -> None:
    font = _preview_font(22)
    line_height = 28
    cursor_y = y
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            cursor_y += line_height // 2
            continue
        for part in _wrap_text(line, width, font):
            if cursor_y + line_height > y + height:
                return
            draw.text((x, cursor_y), part, fill="#101828", font=font)
            cursor_y += line_height


def _wrap_text(text: str, width: int, font) -> list[str]:
    result: list[str] = []
    current = ""
    for char in text:
        candidate = current + char
        try:
            size = font.getbbox(candidate)[2]
        except Exception:
            size = len(candidate) * 12
        if size > width and current:
            result.append(current)
            current = char
        else:
            current = candidate
    if current:
        result.append(current)
    return result


def _preview_font(size: int):
    for font_name in ("msyh.ttc", "simhei.ttf", "arial.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except Exception:
            continue
    return ImageFont.load_default()
