from __future__ import annotations

import json
import re
import uuid
from copy import deepcopy
from pathlib import Path
from typing import Any, AsyncGenerator, Optional

from anthropic import Anthropic
from openai import OpenAI

from app.api.routes.ai_ppt import ai_ppt_base_dir, resolve_ai_ppt_file, resolve_ai_ppt_template
from app.config import get_settings
from app.skills.base import BaseSkill, SkillContext, SkillResult

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    Inches = None
    Pt = None

MIN_SLIDES = 3
MAX_SLIDES = 18

PPT_KEYWORDS = ("ppt", "PPT", "幻灯片", "演示文稿", "演示稿", "slides", "slide")
PPT_ACTION_KEYWORDS = ("做", "生成", "创建", "新建", "制作", "改", "修改", "优化", "润色", "补", "删", "调整")


def is_ai_ppt_request(text: str, requested_skill: str = "") -> bool:
    if requested_skill == "ai_ppt":
        return True
    normalized = text or ""
    if not any(keyword in normalized for keyword in PPT_KEYWORDS):
        return False
    return any(keyword in normalized for keyword in PPT_ACTION_KEYWORDS) or "ai" in normalized.lower()


class AIPPTSkill(BaseSkill):
    """Create and revise editable PPTX files, optimized for follow-up Feishu use."""

    def __init__(self) -> None:
        self.settings = get_settings()
        self.output_dir = ai_ppt_base_dir()
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.client: Optional[Any] = None
        if self.settings.LLM_PROVIDER == "anthropic" and self.settings.ANTHROPIC_API_KEY:
            self.client = Anthropic(api_key=self.settings.ANTHROPIC_API_KEY)
        elif self.settings.OPENAI_API_KEY:
            self.client = OpenAI(api_key=self.settings.OPENAI_API_KEY, base_url=self.settings.OPENAI_BASE_URL)

    @property
    def name(self) -> str:
        return "ai_ppt"

    @property
    def description(self) -> str:
        return (
            "飞书 CLI 演示文稿增强能力：根据对话生成可编辑 PPTX，支持基于上一版继续修改、上传源 PPT 后润色、"
            "套用 ppt-master 风格模板，并给出飞书上传/继续编辑入口。"
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "用户关于生成或修改 PPT 的完整需求"},
                "action": {"type": "string", "enum": ["create", "revise"]},
                "slide_count": {"type": "integer", "description": "期望页数，范围 3-18"},
                "style": {"type": "string", "description": "商务、科技、咨询、极简、培训、路演等"},
                "template_id": {"type": "string", "description": "模板 ID"},
                "source_ppt_id": {"type": "string", "description": "上传的待修改 PPT 文件 ID"},
            },
            "required": ["query"],
        }

    async def execute(self, context: SkillContext, **kwargs: Any) -> SkillResult:
        chunks: list[str] = []
        metadata: dict[str, Any] = {}
        async for chunk in self.execute_stream(context, **kwargs):
            if chunk.get("type") == "content":
                chunks.append(str(chunk.get("content") or ""))
            elif chunk.get("type") == "metadata":
                metadata.update(chunk.get("data") or {})
        return SkillResult(success=True, message="".join(chunks), data=metadata or None)

    async def execute_stream(self, context: SkillContext, **kwargs: Any) -> AsyncGenerator[dict[str, Any], None]:
        query = str(kwargs.get("query") or context.message)
        action = self._resolve_action(query, kwargs.get("action"))
        slide_count = self._resolve_slide_count(query, kwargs.get("slide_count"))
        style = str(kwargs.get("style") or self._infer_style(query))
        template_id = kwargs.get("template_id")
        source_ppt_id = kwargs.get("source_ppt_id")

        if Presentation is None:
            message = "飞书演示文稿增强能力需要安装 python-pptx 和 Pillow 后才能生成 PPTX，请更新依赖并重启服务。"
            yield {"type": "content", "content": message}
            yield {"type": "metadata", "data": {"setup_required": True, "missing_dependency": "python-pptx"}}
            return

        yield {"type": "progress", "content": "飞书 CLI：正在理解演示文稿主题、受众、页数和飞书使用场景"}
        previous = self._latest_ai_ppt_from_history(context)
        previous_path = Path(previous["path"]) if previous and previous.get("path") else None
        previous_version = int(previous.get("version", 0)) if previous else 0

        source_path = resolve_ai_ppt_file(source_ppt_id) if source_ppt_id else None
        template_path = resolve_ai_ppt_template(template_id)
        base_path: Optional[Path] = None
        if action == "revise" and previous_path and previous_path.is_file():
            base_path = previous_path
            yield {"type": "progress", "content": f"飞书 CLI：已接上上一版演示文稿 {previous_path.name}"}
        elif source_path:
            base_path = source_path
            action = "revise"
            yield {"type": "progress", "content": f"飞书 CLI：已读取上传源文件 {source_path.name}"}
        elif template_path:
            base_path = template_path
            yield {"type": "progress", "content": f"飞书 CLI：已套用演示文稿模板 {template_path.name}"}

        source_outline = previous.get("outline") if action == "revise" and previous else None
        yield {"type": "progress", "content": "飞书 CLI：正在生成结构化大纲和讲述顺序"}
        outline = await self._build_outline(query, slide_count, style, source_outline)

        yield {"type": "progress", "content": "飞书 CLI：正在排版为可编辑 PPTX"}
        result = self._render_deck(
            outline=outline,
            query=query,
            style=style,
            base_path=base_path,
            template_id=template_id,
            source_ppt_id=source_ppt_id,
            previous=previous,
            version=previous_version + 1 if previous else 1,
        )

        yield {"type": "progress", "content": "飞书 CLI：文件已生成，正在准备预览和飞书继续编辑建议"}
        message = self._build_response_message(result, action, style)
        for piece in re.findall(r".{1,36}", message, flags=re.S):
            yield {"type": "content", "content": piece}

        trace = [
            "飞书 CLI：演示文稿需求解析完成",
            f"飞书 CLI：生成 {len(result['slides'])} 页大纲",
            f"飞书 CLI：输出文件 {result['filename']}",
        ]
        yield {
            "type": "metadata",
            "data": {
                "ai_ppt": result,
                "last_ai_ppt": result,
                "execution_trace": trace,
                "lark_progress": trace,
            },
        }

    def _latest_ai_ppt_from_history(self, context: SkillContext) -> Optional[dict[str, Any]]:
        if isinstance(context.metadata, dict) and isinstance(context.metadata.get("last_ai_ppt"), dict):
            return context.metadata["last_ai_ppt"]
        for item in reversed(context.history or []):
            metadata = item.get("metadata") if isinstance(item, dict) else None
            if isinstance(metadata, dict) and isinstance(metadata.get("last_ai_ppt"), dict):
                return metadata["last_ai_ppt"]
            if isinstance(metadata, dict) and isinstance(metadata.get("ai_ppt"), dict):
                return metadata["ai_ppt"]
        return None

    def _needs_feishu_followup(self, query: str) -> bool:
        text = (query or "").lower()
        return any(
            keyword in text
            for keyword in (
                "发给",
                "发送",
                "转发",
                "群",
                "群里",
                "私聊",
                "上传",
                "保存",
                "云空间",
                "云文档",
                "drive",
                "chat",
                "send",
                "share",
                "upload",
            )
        )

    def _build_followup_plan(self, result: dict[str, Any], query: str) -> dict[str, Any]:
        actions = []
        if self._asks_upload(query):
            actions.append(
                {
                    "command": f"lark-cli drive +upload --file {result['path']} --as user",
                    "reason": "将刚生成的 PPTX 上传到飞书云空间或用户指定的云文档位置。",
                    "expected": "write",
                }
            )
        if self._asks_send(query):
            actions.append(
                {
                    "command": f"lark-cli im +messages-send --file {result['path']} --as user",
                    "reason": "按用户指定的联系人或群聊发送刚生成的 PPTX；执行前会先搜索并解析目标。",
                    "expected": "write",
                }
            )
        return {
            "summary": "生成 PPTX 后继续执行飞书 CLI 上传/发送动作。",
            "intent_type": "lark_slides_ppt_followup",
            "relevant_skills": ["lark-drive", "lark-im", "lark-contact", "lark-shared"],
            "references": ["lark-drive-upload.md", "lark-im-messages-send.md", "lark-im-chat-search.md"],
            "need_confirmation": True,
            "reason_for_confirmation": "该请求会上传或发送文件到飞书，属于写操作。",
            "commands": actions,
        }

    def _build_followup_confirmation_message(self, result: dict[str, Any], query: str) -> str:
        plan = self._build_followup_plan(result, query)
        commands = "\n".join(f"- `{item['command']}`：{item['reason']}" for item in plan["commands"])
        return (
            f"飞书演示文稿已生成：**{result['title']}**。\n\n"
            "我检测到你还要求继续上传、保存、转发或发送到飞书联系人/群聊。"
            "这些是飞书写操作，需要确认后继续执行。\n\n"
            f"{commands}\n\n"
            "确认后我会继续用飞书 CLI 完成后续动作。"
        )

    def _asks_upload(self, query: str) -> bool:
        return any(keyword in query for keyword in ("上传", "保存", "云空间", "云文档", "drive", "upload"))

    def _asks_send(self, query: str) -> bool:
        return any(keyword in query for keyword in ("发给", "发送", "转发", "群", "群里", "私聊", "chat", "send", "share"))

    async def _execute_feishu_followup(
        self,
        context: SkillContext,
        query: str,
        result: dict[str, Any],
        timeout: int,
    ) -> AsyncGenerator[dict[str, Any], None]:
        from app.skills.lark_cli.skill import LarkCLISkill

        followup_query = self._build_followup_query(query, result)
        skill = LarkCLISkill()
        async for event in skill.execute_stream(
            context,
            query=followup_query,
            confirm_write=True,
            timeout=timeout,
        ):
            if event.get("type") == "progress":
                content = str(event.get("content") or "")
                yield {"type": "progress", "content": f"飞书 CLI 后续动作：{content}"}
            else:
                yield event

    def _build_followup_query(self, query: str, result: dict[str, Any]) -> str:
        path = result["path"]
        name = result["filename"]
        action_hints: list[str] = []
        if self._asks_upload(query):
            action_hints.append(
                "如果用户要求上传/保存到云空间或云文档，使用 `lark-cli drive +upload --file <本地PPT路径> --as user`；"
                "若用户给出 folder-token 或 wiki-token，按 drive +upload 参数传入。"
            )
        if self._asks_send(query):
            action_hints.append(
                "如果用户要求发给某人，先用联系人搜索拿 open_id，再用 `lark-cli im +messages-send --user-id <open_id> --file <本地PPT路径> --as user`；"
                "如果用户要求发到群里，先搜索群拿 chat_id，再用 `lark-cli im +messages-send --chat-id <chat_id> --file <本地PPT路径> --as user`。"
            )
        return (
            "这是一个飞书 CLI 演示文稿后续动作任务。刚刚已经生成了可发送/上传的 PPTX 文件。\n"
            f"原始用户需求：{query}\n"
            f"本地 PPTX 路径：{path}\n"
            f"PPTX 文件名：{name}\n"
            "请严格围绕原始需求继续闭环执行飞书动作，不要重新生成 PPT。\n"
            + "\n".join(action_hints)
        )

    def _resolve_action(self, query: str, action: Optional[str]) -> str:
        if action in {"create", "revise"}:
            return action
        revise_words = ("修改", "改一下", "调整", "润色", "替换", "删掉", "加一页", "补充", "上一版", "上一份")
        return "revise" if any(word in query for word in revise_words) else "create"

    def _resolve_slide_count(self, query: str, value: Any) -> int:
        if isinstance(value, int):
            return max(MIN_SLIDES, min(MAX_SLIDES, value))
        match = re.search(r"(\d{1,2})\s*(?:页|张|p|P|slide|slides)", query)
        if match:
            return max(MIN_SLIDES, min(MAX_SLIDES, int(match.group(1))))
        if any(word in query for word in ("路演", "方案", "汇报", "商业计划")):
            return 8
        return 6

    def _infer_style(self, query: str) -> str:
        for style in ("科技", "商务", "咨询", "极简", "培训", "路演", "学术", "活泼"):
            if style in query:
                return style
        return "商务简洁"

    async def _build_outline(
        self,
        query: str,
        slide_count: int,
        style: str,
        source_outline: Optional[dict[str, Any]] = None,
    ) -> dict[str, Any]:
        try:
            raw = await self._call_llm(self._outline_prompt(query, slide_count, style, source_outline))
            return self._normalize_outline(self._extract_json(raw), query, slide_count)
        except Exception:
            return self._fallback_outline(query, slide_count, style, source_outline)

    def _outline_prompt(
        self,
        query: str,
        slide_count: int,
        style: str,
        source_outline: Optional[dict[str, Any]],
    ) -> str:
        previous = ""
        if source_outline:
            previous = "\n上一版 PPT 大纲 JSON：\n" + json.dumps(source_outline, ensure_ascii=False)
        return f"""你是资深演示文稿策划师，请把用户需求转成可编辑 PPT 大纲。
要求：
- 只输出 JSON，不要 Markdown。
- slides 必须正好 {slide_count} 页。
- 每页包含 title、bullets、speaker_notes、layout。
- bullets 每页 3-5 条，每条不超过 32 个汉字。
- layout 可选 cover、section、bullets、two_column、timeline、closing。
- 内容要能直接落到 PowerPoint 和飞书 Slides 中，不要出现“待补充”等空话。
- 标题短、信息密度高，适合在飞书里继续协作修改。

视觉风格：{style}
用户需求：{query}
{previous}

JSON 结构：
{{
  "title": "PPT 标题",
  "subtitle": "一句副标题",
  "slides": [
    {{
      "title": "页标题",
      "bullets": ["要点1", "要点2"],
      "speaker_notes": "给演讲者的简短提示",
      "layout": "bullets"
    }}
  ]
}}"""

    async def _call_llm(self, prompt: str) -> str:
        if not self.client:
            raise RuntimeError("LLM client is not configured")
        if self.settings.LLM_PROVIDER == "anthropic":
            response = self.client.messages.create(
                model=self.settings.LLM_MODEL,
                max_tokens=4096,
                system="你只输出可解析 JSON。",
                messages=[{"role": "user", "content": prompt}],
            )
            return "".join(block.text for block in response.content if block.type == "text")
        response = self.client.chat.completions.create(
            model=self.settings.LLM_MODEL,
            messages=[
                {"role": "system", "content": "你只输出可解析 JSON。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.35,
        )
        return response.choices[0].message.content or ""

    def _extract_json(self, raw: str) -> dict[str, Any]:
        text = raw.strip()
        if text.startswith("```"):
            text = re.sub(r"^```(?:json)?", "", text).strip()
            text = re.sub(r"```$", "", text).strip()
        start = text.find("{")
        end = text.rfind("}")
        if start >= 0 and end > start:
            text = text[start : end + 1]
        return json.loads(text)

    def _normalize_outline(self, outline: dict[str, Any], query: str, slide_count: int) -> dict[str, Any]:
        title = str(outline.get("title") or self._title_from_query(query)).strip()
        raw_slides = outline.get("slides") if isinstance(outline.get("slides"), list) else []
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(raw_slides[:slide_count], start=1):
            if not isinstance(slide, dict):
                continue
            bullets = slide.get("bullets") if isinstance(slide.get("bullets"), list) else []
            slides.append(
                {
                    "title": str(slide.get("title") or f"第 {index} 页").strip(),
                    "bullets": [str(item).strip() for item in bullets if str(item).strip()][:5],
                    "speaker_notes": str(slide.get("speaker_notes") or "").strip(),
                    "layout": str(slide.get("layout") or "bullets").strip(),
                }
            )
        while len(slides) < slide_count:
            slides.append(
                {
                    "title": f"关键内容 {len(slides) + 1}",
                    "bullets": ["核心观点", "支撑事实", "落地建议"],
                    "speaker_notes": "",
                    "layout": "bullets",
                }
            )
        slides[0]["layout"] = "cover"
        slides[-1]["layout"] = "closing"
        return {"title": title, "subtitle": str(outline.get("subtitle") or "").strip(), "slides": slides}

    def _fallback_outline(
        self,
        query: str,
        slide_count: int,
        style: str,
        source_outline: Optional[dict[str, Any]],
    ) -> dict[str, Any]:
        if source_outline:
            outline = deepcopy(source_outline)
            outline["subtitle"] = f"已根据新要求调整：{query[:40]}"
            slides = outline.get("slides", [])
            if slides:
                slides[-1]["bullets"] = ["聚焦本次修改要求", "优化叙事顺序与表达", "补强结论和行动建议"]
            return self._normalize_outline(outline, query, min(MAX_SLIDES, max(MIN_SLIDES, len(slides))))

        title = self._title_from_query(query)
        sections = [
            ("背景与目标", ["问题背景清晰化", "明确本次汇报目标", "对齐受众关注点"]),
            ("核心洞察", ["提炼关键事实", "呈现主要判断", "突出机会与风险"]),
            ("方案设计", ["形成总体路径", "拆解关键举措", "匹配资源与节奏"]),
            ("实施计划", ["划分阶段任务", "明确责任分工", "设置里程碑检查"]),
            ("价值预期", ["说明业务收益", "量化关键指标", "建立复盘机制"]),
            ("下一步行动", ["确认决策事项", "启动近期动作", "持续跟踪反馈"]),
        ]
        slides: list[dict[str, Any]] = []
        for index in range(slide_count):
            if index == 0:
                slides.append(
                    {
                        "title": title,
                        "bullets": [f"风格：{style}", "面向决策与执行", "结构清晰、可直接汇报"],
                        "speaker_notes": "开场说明汇报目的和预期产出。",
                        "layout": "cover",
                    }
                )
                continue
            name, bullets = sections[(index - 1) % len(sections)]
            slides.append(
                {
                    "title": name,
                    "bullets": bullets,
                    "speaker_notes": "围绕事实、判断、动作三段展开。",
                    "layout": "closing" if index == slide_count - 1 else "bullets",
                }
            )
        return {"title": title, "subtitle": "飞书 CLI 自动生成", "slides": slides}

    def _title_from_query(self, query: str) -> str:
        cleaned = re.sub(r"(帮我|做一份|生成|创建|制作|PPT|ppt|幻灯片|演示文稿|关于)", " ", query)
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" ：:，,。")
        return cleaned[:28] or "AI 生成演示文稿"

    def _render_deck(
        self,
        outline: dict[str, Any],
        query: str,
        style: str,
        base_path: Optional[Path] = None,
        template_id: Optional[str] = None,
        source_ppt_id: Optional[str] = None,
        previous: Optional[dict[str, Any]] = None,
        version: int = 1,
    ) -> dict[str, Any]:
        prs = Presentation(str(base_path)) if base_path and base_path.is_file() else Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        self._delete_all_slides(prs)

        palette = self._palette(style)
        for index, slide_data in enumerate(outline["slides"], start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_background(slide, palette["bg"])
            layout = slide_data.get("layout")
            if index == 1 or layout == "cover":
                self._draw_cover(slide, outline, slide_data, palette)
            elif layout == "two_column":
                self._draw_two_column(slide, slide_data, palette, index, len(outline["slides"]))
            else:
                self._draw_bullets(slide, slide_data, palette, index, len(outline["slides"]))
            if slide_data.get("speaker_notes"):
                slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

        safe_title = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", outline["title"]).strip("_")[:32] or "ai_ppt"
        project_id = (previous or {}).get("project_id") or uuid.uuid4().hex
        filename = f"{project_id}_v{version:03d}_{safe_title}.pptx"
        output_path = self.output_dir / filename
        prs.save(output_path)

        return {
            "title": outline["title"],
            "subtitle": outline.get("subtitle") or "",
            "filename": filename,
            "download_url": f"/api/v1/ai-ppt/files/{filename}",
            "preview_url": f"/api/v1/ai-ppt/files/{filename}/preview",
            "path": str(output_path),
            "query": query,
            "style": style,
            "template_id": template_id,
            "source_ppt_id": source_ppt_id,
            "source_filename": str(base_path) if base_path else None,
            "project_id": project_id,
            "version": version,
            "previous_filename": (previous or {}).get("filename"),
            "outline": outline,
            "slides": [
                {"index": idx + 1, "title": slide["title"], "layout": slide.get("layout", "bullets")}
                for idx, slide in enumerate(outline["slides"])
            ],
            "feishu_tip": "下载 PPTX 后可上传到飞书云空间或用飞书 Slides 打开继续协作；也可以继续在当前飞书 CLI 对话里说修改第几页。",
        }

    def _delete_all_slides(self, prs: Any) -> None:
        slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001
        for slide_id in slide_ids:
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001

    def _palette(self, style: str) -> dict[str, str]:
        if "科技" in style:
            return {"bg": "F7FBFF", "fg": "172033", "accent": "1167B1", "muted": "60758D", "soft": "EAF4FF"}
        if "咨询" in style:
            return {"bg": "FBFCFD", "fg": "111827", "accent": "0F172A", "muted": "64748B", "soft": "EEF2F7"}
        if "活泼" in style:
            return {"bg": "FFF8F1", "fg": "252525", "accent": "E85D3F", "muted": "76685F", "soft": "FFE9DF"}
        if "极简" in style:
            return {"bg": "FFFFFF", "fg": "161616", "accent": "111827", "muted": "6B7280", "soft": "F3F4F6"}
        return {"bg": "FAFBFC", "fg": "18202A", "accent": "2454D6", "muted": "667085", "soft": "EEF3FF"}

    def _set_background(self, slide: Any, color: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color)

    def _draw_cover(self, slide: Any, outline: dict[str, Any], slide_data: dict[str, Any], palette: dict[str, str]) -> None:
        self._add_shape(slide, 0, 0, 0.18, 7.5, palette["accent"])
        self._add_shape(slide, 9.65, 0, 3.7, 7.5, palette["soft"])
        self._add_text(slide, outline["title"], 0.9, 1.6, 10.4, 1.0, 36, palette["fg"], bold=True)
        subtitle = outline.get("subtitle") or "由 Feishu CLI Web 演示文稿能力生成"
        self._add_text(slide, subtitle, 0.95, 2.75, 9.8, 0.45, 16, palette["muted"])
        self._add_text(slide, "Feishu CLI / Slides-ready", 0.95, 5.9, 3.2, 0.35, 13, palette["accent"], bold=True)
        for idx, item in enumerate(slide_data.get("bullets", [])[:3]):
            self._add_chip(slide, item, 0.95 + idx * 2.8, 6.34, 2.55, palette)

    def _draw_bullets(self, slide: Any, slide_data: dict[str, Any], palette: dict[str, str], index: int, total: int) -> None:
        self._add_header(slide, slide_data["title"], palette, index)
        for bullet_index, bullet in enumerate(slide_data.get("bullets", [])[:5], start=1):
            y = 1.72 + (bullet_index - 1) * 0.78
            self._add_number(slide, bullet_index, 1.0, y + 0.03, palette)
            self._add_text(slide, bullet, 1.55, y, 10.5, 0.42, 21, palette["fg"], bold=bullet_index == 1)
        self._add_footer(slide, palette, index, total)

    def _draw_two_column(self, slide: Any, slide_data: dict[str, Any], palette: dict[str, str], index: int, total: int) -> None:
        self._add_header(slide, slide_data["title"], palette, index)
        bullets = slide_data.get("bullets", [])[:6]
        midpoint = max(1, (len(bullets) + 1) // 2)
        for col, items in enumerate((bullets[:midpoint], bullets[midpoint:])):
            x = 0.95 + col * 5.95
            self._add_shape(slide, x, 1.62, 5.25, 4.7, "FFFFFF")
            self._add_text(slide, "重点" if col == 0 else "行动", x + 0.35, 1.92, 4.4, 0.35, 15, palette["accent"], bold=True)
            for idx, bullet in enumerate(items):
                self._add_text(slide, f"- {bullet}", x + 0.35, 2.42 + idx * 0.62, 4.55, 0.35, 17, palette["fg"])
        self._add_footer(slide, palette, index, total)

    def _add_header(self, slide: Any, title: str, palette: dict[str, str], index: int) -> None:
        self._add_text(slide, title, 0.85, 0.52, 10.6, 0.55, 26, palette["fg"], bold=True)
        self._add_text(slide, f"{index:02d}", 11.8, 0.55, 0.75, 0.35, 14, palette["accent"], bold=True, align="right")
        self._add_shape(slide, 0.86, 1.2, 1.25, 0.04, palette["accent"])

    def _add_footer(self, slide: Any, palette: dict[str, str], index: int, total: int) -> None:
        self._add_text(slide, "Feishu CLI Web / Slides", 0.9, 6.92, 2.8, 0.25, 9, palette["muted"])
        self._add_text(slide, f"{index}/{total}", 11.6, 6.92, 0.85, 0.25, 9, palette["muted"], align="right")

    def _add_chip(self, slide: Any, text: str, x: float, y: float, width: float, palette: dict[str, str]) -> None:
        shape = self._add_shape(slide, x, y, width, 0.42, "FFFFFF")
        shape.line.color.rgb = RGBColor.from_string("D9E2F2")
        self._add_text(slide, text, x + 0.15, y + 0.08, width - 0.3, 0.22, 9, palette["muted"])

    def _add_number(self, slide: Any, number: int, x: float, y: float, palette: dict[str, str]) -> None:
        shape = self._add_shape(slide, x, y, 0.34, 0.34, palette["accent"])
        shape.line.color.rgb = RGBColor.from_string(palette["accent"])
        self._add_text(slide, str(number), x, y + 0.04, 0.34, 0.18, 9, "FFFFFF", bold=True, align="center")

    def _add_shape(self, slide: Any, x: float, y: float, width: float, height: float, color: str) -> Any:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        shape.line.color.rgb = RGBColor.from_string(color)
        return shape

    def _add_text(
        self,
        slide: Any,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        size: int,
        color: str,
        bold: bool = False,
        align: str = "left",
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        return box

    def _build_response_message(self, result: dict[str, Any], action: str, style: str, followup_text: str = "") -> str:
        verb = "修改好了" if action == "revise" else "做好了"
        slide_titles = "\n".join(f"{slide['index']}. {slide['title']}" for slide in result["slides"])
        version_text = f"v{int(result.get('version', 1)):03d}"
        previous = result.get("previous_filename")
        source_line = f"- 基于：{previous}\n" if previous else ""
        followup_line = f"\n飞书后续动作结果：\n{followup_text.strip()}\n" if followup_text.strip() else ""
        return (
            f"飞书演示文稿 {verb}：**{result['title']}**（{version_text}）\n\n"
            f"- 页数：{len(result['slides'])} 页\n"
            f"- 风格：{style}\n"
            f"{source_line}"
            f"- 文件：[{result['filename']}]({result['download_url']})\n\n"
            f"页面结构：\n{slide_titles}\n\n"
            f"{followup_line}"
            "你可以继续直接说“把第 3 页改得更有说服力”或“再加一页实施计划”，我会基于这份继续调整。"
        )
